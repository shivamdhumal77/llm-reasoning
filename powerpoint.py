import os
from pptx import Presentation
import win32com.client

# Function to extract text from a .pptx file
def extract_text_from_pptx(pptx_file_path):
    presentation = Presentation(pptx_file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from a .ppt file (requires pywin32 on Windows)
def extract_text_from_ppt(ppt_file_path):
    # Initialize PowerPoint application
    power_point = win32com.client.Dispatch("PowerPoint.Application")
    presentation = power_point.Presentations.Open(ppt_file_path)
    text = ""
    for slide in presentation.Slides:
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                if shape.TextFrame.HasText:
                    text += shape.TextFrame.TextRange.Text + "\n"
    presentation.Close()
    power_point.Quit()
    return text

# Function to extract text from both .pptx and .ppt files
def extract_text_from_ppt_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_extension == '.ppt':
        return extract_text_from_ppt(file_path)
    else:
        raise ValueError("Unsupported file type. Please upload a .ppt or .pptx file.")
